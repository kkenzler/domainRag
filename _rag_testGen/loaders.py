from __future__ import annotations

"""loaders.py — Document text extraction.

Responsibilities:
  - Load raw text from PDF, PPTX, DOCX, TXT/MD files
  - Light structural preprocessing (line wrap repair, artifact removal)
  - Return LoadedDoc with stable sha256 identity

Chunking is NOT done here. The context model in ingest.py
receives the full document text and produces knowledge chunks.
"""

import hashlib
import re
import sys

from dataclasses import dataclass
from pathlib import Path
from typing import Optional


@dataclass(frozen=True)
class LoadedDoc:
    """A loaded document with stable identity for traceability and idempotent upserts."""
    path: Path
    sha256: str
    text: str          # full preprocessed text of the document
    page_count: int    # number of pages/slides (for logging)


def sha256_file(path: Path) -> str:
    """Computes sha256 of a file for stable document identity."""
    h = hashlib.sha256()
    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(1024 * 1024), b""):
            h.update(chunk)
    return h.hexdigest()


def load_text_file(path: Path) -> tuple[str, int]:
    """Returns (text, page_count=1) for plain text files."""
    return path.read_text(encoding="utf-8", errors="replace"), 1


def load_docx(path: Path) -> tuple[str, int]:
    """Loads a .docx file using python-docx; returns (text, paragraph_count)."""
    from docx import Document
    doc = Document(str(path))
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(parts).strip(), len(parts)


def load_pdf(path: Path) -> tuple[str, int]:
    """Loads a PDF using PyMuPDF (preferred) or pypdf fallback.
    Returns (text, page_count).
    """
    # Try PyMuPDF first
    try:
        import fitz  # type: ignore
        doc = fitz.open(str(path))
        parts: list[str] = []
        for i in range(doc.page_count):
            page = doc.load_page(i)
            t = (page.get_text("text") or "").strip()
            if t:
                parts.append(f"--- Page {i + 1} ---\n{t}")
        return "\n\n".join(parts).strip(), doc.page_count
    except ImportError:
        pass
    except Exception as e:
        print(f"  [WARNING] PyMuPDF failed on {path.name}: {e}", file=sys.stderr)

    # Fallback: pypdf
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        raise RuntimeError(
            "PDF support requires PyMuPDF (fitz) or pypdf. "
            "Install one: pip install pymupdf  OR  pip install pypdf"
        )
    reader = PdfReader(str(path))
    parts = []
    for i, page in enumerate(reader.pages, start=1):
        t = (page.extract_text() or "").strip()
        if t:
            parts.append(f"--- Page {i} ---\n{t}")
    return "\n\n".join(parts).strip(), len(reader.pages)


def load_pptx(path: Path) -> tuple[str, int]:
    """Loads a PPTX using python-pptx; extracts text, tables, and speaker notes.
    Returns (text, slide_count).
    """
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise RuntimeError(
            "PPTX support requires python-pptx. Install: pip install python-pptx"
        )
    prs = Presentation(str(path))
    parts: list[str] = []

    for si, slide in enumerate(prs.slides, start=1):
        slide_lines: list[str] = [f"--- Slide {si} ---"]

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                txt = (getattr(shape, "text") or "").strip()
                if txt:
                    slide_lines.append(txt)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        txt = (cell.text or "").strip()
                        if txt:
                            slide_lines.append(txt)

        if getattr(slide, "has_notes_slide", False):
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                slide_lines.append(f"[Notes] {notes_text}")

        if len(slide_lines) > 1:
            parts.append("\n".join(slide_lines))

    return "\n\n".join(parts).strip(), len(prs.slides)


# ---------------------------------------------------------------------------
# Preprocessing
# ---------------------------------------------------------------------------

_TS_PAT = re.compile(r"^\[?\d{1,2}:\d{2}(?::\d{2})?(?:\.\d{1,3})?\]?$")
_INT_PAT = re.compile(r"^\d{1,3}$")
_WS_RE = re.compile(r"[ \t]+")


def _unwrap_pdf_lines(lines: list[str]) -> list[str]:
    """Heuristically joins hard-wrapped PDF lines while preserving blank lines."""
    out: list[str] = []
    i = 0
    while i < len(lines):
        cur = lines[i]
        if cur == "":
            out.append("")
            i += 1
            continue
        merged = cur
        j = i + 1
        while j < len(lines):
            nxt = lines[j]
            if nxt == "":
                break
            if re.search(r"[.!?]\s*$", merged):
                break
            if re.match(r"^[a-z(]", nxt):
                merged = merged.rstrip() + " " + nxt.lstrip()
                j += 1
                continue
            break
        out.append(merged)
        i = j
    return out


def preprocess_text(text: str, source_ext: str = "") -> str:
    """Normalizes extracted text while preserving paragraph boundaries."""
    t = (text or "")
    t = t.replace("\r\n", "\n").replace("\r", "\n")
    # De-hyphenate PDF line wrap artifacts
    t = re.sub(r"(\w)-\n(\w)", r"\1\2", t)
    t = _WS_RE.sub(" ", t)

    raw_lines = t.split("\n")
    lines = [ln.strip() for ln in raw_lines]

    # Drop timestamp-only lines (transcript artifacts)
    lines = [("" if _TS_PAT.match(ln) else ln) for ln in lines]
    # Drop standalone page/slide number artifacts
    lines = [("" if (ln and _INT_PAT.match(ln)) else ln) for ln in lines]
    # Drop Wingdings 'z' bullet artifacts
    lines = [("" if ln == "z" else ln) for ln in lines]
    lines = [re.sub(r"^\s*z\s+", "", ln) for ln in lines]

    # Remove repeated header/footer lines (freq >= 4, length <= 120)
    freq: dict[str, int] = {}
    for ln in lines:
        if ln:
            freq[ln.lower()] = freq.get(ln.lower(), 0) + 1
    cleaned: list[str] = []
    for ln in lines:
        if ln and freq.get(ln.lower(), 0) >= 4 and len(ln) <= 120:
            cleaned.append("")
        else:
            cleaned.append(ln)

    if source_ext.lower() == ".pdf":
        cleaned = _unwrap_pdf_lines(cleaned)

    out = "\n".join(cleaned)
    out = re.sub(r"\n{3,}", "\n\n", out)
    return out.strip()


def _is_spoken_math_transcript(text: str) -> bool:
    """Detects auto-generated spoken-math transcripts (unusable for generation)."""
    markers = [
        "superscript", "subscript", "open parenthesis", "close parenthesis",
        "divided by", "square root of", "equals sign", "times sign",
    ]
    hits = sum(1 for m in markers if m in text.lower())
    return hits >= 3


# ---------------------------------------------------------------------------
# MP4 — auto-transcribe via faster-whisper, sidecar .txt cache
# ---------------------------------------------------------------------------

def load_mp4(path: Path) -> tuple[str, int]:
    """Transcribes an MP4 using faster-whisper (local, no API).

    Checks for a sidecar .txt (same stem, same directory) first.
    If found, loads it directly — avoids re-transcribing on subsequent ingest runs.
    If not found, transcribes and writes the sidecar, then returns the text.
    """
    sidecar = path.with_suffix(".txt")

    if sidecar.exists():
        print(f"  [mp4] Using existing transcript: {sidecar.name}", flush=True)
        text = sidecar.read_text(encoding="utf-8", errors="replace")
        return text, 1

    print(f"  [mp4] No transcript found — transcribing: {path.name}", flush=True)

    try:
        from faster_whisper import WhisperModel
    except ImportError:
        raise RuntimeError(
            "MP4 transcription requires faster-whisper. "
            "Install: pip install faster-whisper"
        )

    try:
        model = WhisperModel("small", device="cuda", compute_type="float16")
        print("  [mp4] Using GPU.", flush=True)
    except Exception:
        model = WhisperModel("small", device="cpu", compute_type="int8")
        print("  [mp4] Using CPU (CUDA unavailable).", flush=True)

    segments, info = model.transcribe(str(path), language="en")
    duration = info.duration

    lines = []
    for segment in segments:
        h = int(segment.start // 3600)
        m = int((segment.start % 3600) // 60)
        s = int(segment.start % 60)
        lines.append(f"[{h:02}:{m:02}:{s:02}] {segment.text.strip()}")
        pct = min(segment.end / duration, 1.0) if duration > 0 else 0
        filled = int(40 * pct)
        bar = "#" * filled + "-" * (40 - filled)
        sys.stderr.write(f"\r  [mp4] [{bar}] {pct*100:.1f}%")
        sys.stderr.flush()

    sys.stderr.write("\n")
    sys.stderr.flush()

    # Free CUDA resources before returning — prevents stack overrun on Windows
    # when the next pipeline stage initialises its own native libraries.
    del model
    try:
        import torch
        torch.cuda.empty_cache()
    except Exception:
        pass

    print(f"  [mp4] Done. Writing transcript: {sidecar.name}", flush=True)
    text = "\n".join(lines)
    sidecar.write_text(text, encoding="utf-8")

    return text, 1


# ---------------------------------------------------------------------------
# Public interface
# ---------------------------------------------------------------------------

def load_document(path: Path) -> Optional[LoadedDoc]:
    """Loads supported document types; returns None for unsupported or empty docs."""
    path = Path(path)
    ext = path.suffix.lower()
    sha = sha256_file(path)

    if ext in {".txt", ".md"}:
        if path.with_suffix(".mp4").exists():
            return None  # sidecar transcript — handled via the .mp4 entry
        raw_text, page_count = load_text_file(path)
    elif ext == ".docx":
        raw_text, page_count = load_docx(path)
    elif ext == ".pdf":
        raw_text, page_count = load_pdf(path)
    elif ext == ".pptx":
        raw_text, page_count = load_pptx(path)
    elif ext == ".mp4":
        raw_text, page_count = load_mp4(path)
    else:
        return None

    if not raw_text:
        return None

    text = preprocess_text(raw_text, source_ext=ext).strip()
    if not text:
        return None

    if _is_spoken_math_transcript(text):
        print(
            f"  [TRANSCRIPT FLAG] Spoken-math notation detected, skipping: {path.name}",
            file=sys.stderr, flush=True,
        )
        return None

    return LoadedDoc(path=path, sha256=sha, text=text, page_count=page_count)
